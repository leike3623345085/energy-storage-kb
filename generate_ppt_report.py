from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# 创建演示文稿
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle):
    """添加标题页"""
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)
    
    # 添加标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 112, 192)
    p.alignment = PP_ALIGN.CENTER
    
    # 添加副标题
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(100, 100, 100)
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, content_list, two_column=False):
    """添加内容页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 112, 192)
    
    # 内容
    if not two_column:
        content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(12), Inches(5.5))
        tf = content_box.text_frame
        tf.word_wrap = True
        
        for i, item in enumerate(content_list):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(20)
            p.space_after = Pt(12)
    else:
        # 左列
        left_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(5.8), Inches(5.5))
        tf = left_box.text_frame
        tf.word_wrap = True
        mid = len(content_list) // 2
        for i, item in enumerate(content_list[:mid]):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(18)
            p.space_after = Pt(10)
        
        # 右列
        right_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.4), Inches(5.8), Inches(5.5))
        tf = right_box.text_frame
        tf.word_wrap = True
        for i, item in enumerate(content_list[mid:]):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(18)
            p.space_after = Pt(10)
    
    return slide

def add_table_slide(prs, title, headers, rows):
    """添加表格页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 112, 192)
    
    # 表格
    num_rows = len(rows) + 1
    num_cols = len(headers)
    table = slide.shapes.add_table(num_rows, num_cols, Inches(0.7), Inches(1.4), Inches(12), Inches(5.5)).table
    
    # 表头
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(16)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0, 112, 192)
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    # 数据
    for i, row in enumerate(rows):
        for j, cell_text in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(cell_text)
            cell.text_frame.paragraphs[0].font.size = Pt(14)
    
    return slide

# ==================== 开始创建PPT ====================

# 第1页：标题页
add_title_slide(prs, 
    "系统体系总结报告", 
    "储能行业监控系统 v2.1\n报告时间：2026年3月22日")

# 第2页：系统架构总览
add_content_slide(prs, "一、系统架构总览", [
    "三层架构：认知层 → 执行层 → 存储层",
    "认知层（记忆系统）：六大控制面板框架，刚完成重构",
    "执行层（Harness架构）：人类掌舵，智能体执行",
    "存储层（多平台备份）：GitHub + 飞书 + IMA 三重备份",
    "自愈监控系统：闭环修复，规则学习"
])

# 第3页：记忆系统 v2.0
add_table_slide(prs, "二、记忆系统 v2.0 - 六大控制面板", 
    ["文件", "职责", "类比"],
    [
        ["SOUL.md", "人格、风格、原则、边界", "性格DNA"],
        ["USER.md", "用户长期目标、沟通偏好", "用户画像"],
        ["MEMORY.md", "跨项目通用规则、决策", "长期记忆"],
        ["Daily", "当日事件、决定、待办", "日记本"],
        ["TOOLS.md", "执行硬规则、环境配置", "工具箱"],
        ["AGENTS.md", "治理制度、风险策略", "操作手册"]
    ]
)

# 第4页：Harness Engineering 架构
add_content_slide(prs, "三、Harness Engineering 架构", [
    '"人类掌舵，智能体执行" —— OpenAI 2026年2月提出的新范式',
    "Guardrails（护栏）：硬性约束，飞行前检查",
    "Feedback Loop（反馈循环）：错误分类→自动修复→规则学习",
    "Progressive Disclosure（渐进式披露）：分层加载上下文",
    "Agent Runner（执行器）：总控协调，标准化入口"
], two_column=True)

# 第5页：研报分析框架
add_content_slide(prs, "四、研报分析框架（Skill）", [
    "MECE分析法：热点分类，确保不重叠不遗漏",
    "SWOT分析：快速生成行业态势速览",
    "产业链分析：上游→中游→下游动态追踪",
    "影响评估矩阵：事件优先级排序",
    "当前已适配：储能行业，可扩展至其他行业"
])

# 第6页：储能监控系统
add_table_slide(prs, "五、储能监控系统 - 监控覆盖领域",
    ["领域", "关键技术/方向"],
    [
        ["储能系统集成", "压缩空气、飞轮、储热、氢能"],
        ["电池多元化", "锂电、钠电、液流、固态电池"],
        ["PCS与电力电子", "储能变流器、构网型储能"],
        ["消防与安全", "pack级消防、热失控防控"],
        ["智能运维", "BMS/EMS、数字孪生、云平台"],
        ["虚拟电厂", "VPP聚合、V2G、微电网"]
    ]
)

# 第7页：数据流
add_content_slide(prs, "六、数据流", [
    "4网站爬虫（每4小时）→ 中国储能网、北极星、OFweek、高工储能",
    "搜索监控（每天6次）→ 00:00/04:00/08:00/12:00/16:00/20:00",
    "股票行情（交易日9:00/15:00）→ 实时技术指标",
    "报告生成（每日18:00）→ 日报 + 深度分析",
    "多平台同步 → GitHub本地、飞书Bitable、IMA笔记"
])

# 第8页：自愈监控系统
add_table_slide(prs, "七、自愈监控系统成效",
    ["指标", "数值"],
    [
        ["累计检测问题", "50+ 次"],
        ["自动修复成功率", "100%"],
        ["平均修复时间", "< 30秒"],
        ["规则学习条目", "10+ 条"],
        ["运行机制", "检测→分类→修复→验证→记录"]
    ]
)

# 第9页：定时任务体系
add_table_slide(prs, "八、定时任务体系",
    ["任务", "频率", "状态"],
    [
        ["网站爬虫", "每4小时", "✅ ok"],
        ["资讯搜索", "每天6次", "✅ ok"],
        ["系统巡检", "每2小时", "✅ ok"],
        ["日报生成", "每天18:05", "✅ ok"],
        ["深度分析", "每天18:10", "✅ ok"],
        ["飞书同步", "每分钟", "✅ running"],
        ["知识库索引", "每天18:25", "⚠️ error"],
        ["IMA同步", "每2分钟", "⚠️ error"]
    ]
)

# 第10页：系统演进
add_table_slide(prs, "九、系统演进路线图",
    ["阶段", "时间", "里程碑"],
    [
        ["v1.0", "2026-03-03", "系统初始化：爬虫 + 邮件"],
        ["v1.5", "2026-03-11", "自愈系统部署，执行准则"],
        ["v2.0", "2026-03-21", "Harness架构，多平台同步"],
        ["v2.1", "2026-03-22", "记忆重构，研报框架"],
        ["v3.0", "规划中", "模式识别、预测性修复"]
    ]
)

# 第11页：核心设计原则
add_content_slide(prs, "十、核心设计原则", [
    "人类掌舵，智能体执行 —— Harness Engineering 核心",
    "遇到问题直接执行，事后通报 —— 执行准则",
    "分层记忆，结构化存储 —— 记忆系统 v2.0",
    "自愈优先，预防为主 —— 系统稳定性",
    "所有工作流复用 Harness 架构 —— 标准化"
])

# 第12页：结束页
add_title_slide(prs, 
    "感谢观看", 
    "每一层都可以独立迭代，但又相互关联形成闭环\n\n储能行业监控系统 v2.1")

# 保存
output_path = "/root/.openclaw/workspace/系统体系总结报告_20260322.pptx"
prs.save(output_path)
print(f"PPT已生成: {output_path}")
